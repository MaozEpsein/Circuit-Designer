"""
One-shot extractor: read the three IQ/PP/*.pptx files and emit a single
extracted.json describing every slide's text content.

Run from repo root:
    python IQ/PP/extract.py
"""
import json
import os
import sys
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
FILES = [
    "LOGIC.pptx",
    "שאלות קוד.pptx",
    "שאלות מעגלים.pptx",
]


def slide_text(slide):
    title = None
    body_chunks = []
    images = 0
    tables = 0
    shapes_other = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            images += 1
            continue
        if shape.has_table:
            tables += 1
            rows = []
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            body_chunks.append("[TABLE]\n" + "\n".join(rows))
            continue
        if not shape.has_text_frame:
            shapes_other += 1
            continue
        txt = shape.text_frame.text.strip()
        if not txt:
            continue
        # Heuristic: the first non-empty text frame in a layout placeholder
        # is the title if the shape is a placeholder of idx 0.
        if title is None and getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.idx == 0:
                    title = txt
                    continue
            except Exception:
                pass
        body_chunks.append(txt)
    return {
        "title": title,
        "body": "\n\n".join(body_chunks),
        "images": images,
        "tables": tables,
        "shapes_other": shapes_other,
    }


SHORT_NAMES = {
    "LOGIC.pptx": "logic",
    "שאלות קוד.pptx": "code",
    "שאלות מעגלים.pptx": "circuits",
}


def save_images(slide, src_short, slide_idx, out_dir):
    saved = []
    sub_idx = 1
    for shape in slide.shapes:
        if shape.shape_type != 13:  # PICTURE
            continue
        try:
            image = shape.image
        except Exception:
            continue
        ext = image.ext
        fname = f"{src_short}_s{slide_idx:02d}_{sub_idx}.{ext}"
        fpath = os.path.join(out_dir, fname)
        with open(fpath, "wb") as f:
            f.write(image.blob)
        saved.append(fname)
        sub_idx += 1
    return saved


def main():
    out = []
    img_dir = os.path.join(HERE, "slides")
    os.makedirs(img_dir, exist_ok=True)
    for fname in FILES:
        path = os.path.join(HERE, fname)
        if not os.path.exists(path):
            print(f"MISSING: {path}", file=sys.stderr)
            continue
        prs = Presentation(path)
        short = SHORT_NAMES.get(fname, "x")
        for idx, slide in enumerate(prs.slides, start=1):
            info = slide_text(slide)
            imgs = save_images(slide, short, idx, img_dir)
            out.append({
                "source": fname,
                "slide": idx,
                "image_files": imgs,
                **info,
            })
        print(f"{fname}: {len(prs.slides)} slides", file=sys.stderr)
    dest = os.path.join(HERE, "extracted.json")
    with open(dest, "w", encoding="utf-8") as f:
        json.dump(out, f, ensure_ascii=False, indent=2)
    print(f"Wrote {dest} ({len(out)} slides total)", file=sys.stderr)
    print(f"Saved images to {img_dir}", file=sys.stderr)


if __name__ == "__main__":
    main()
